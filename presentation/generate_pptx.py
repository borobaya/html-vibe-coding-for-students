"""Generate a PowerPoint version of the AI & the Future presentation."""

import os
import re
import tempfile
from io import BytesIO
from pathlib import Path

from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Paths ──────────────────────────────────────────────────────
HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
VIBE_ASSETS = HERE.parent / "vibe-coding-intro" / "assets"

# Temp dir for converted images
_tmp = tempfile.mkdtemp(prefix="pptx_icons_")

# Launch headless browser once for all SVG conversions
_pw = sync_playwright().start()
_browser = _pw.chromium.launch()

# ── Animation settings ─────────────────────────────────────────
GIF_FRAME_DELAY_MS = 100    # ~10 fps
GIF_MAX_FRAMES = 120        # hard cap on frames per GIF


def _svg_has_animation(svg_path):
    """Return True if the SVG contains SMIL animation elements."""
    text = Path(svg_path).read_text(errors="ignore")
    return bool(re.search(
        r"<animate\b|<animateTransform\b|<animateMotion\b|<set\b", text
    ))


def _parse_time(val):
    """Convert a SMIL time string like '2s' or '300ms' to float seconds."""
    val = val.strip()
    if val.endswith("ms"):
        return float(val[:-2]) / 1000
    if val.endswith("s"):
        return float(val[:-1])
    try:
        return float(val)
    except ValueError:
        return 0


def _svg_animation_duration(svg_path):
    """Best-effort: find the latest end-point (begin + dur) across all animations.

    For SVGs with repeatCount="indefinite", use the longest repeating dur
    so the GIF captures one full looping cycle.
    """
    text = Path(svg_path).read_text(errors="ignore")
    max_end = 0
    max_repeating_dur = 0
    for m in re.finditer(
        r'<(?:animate|animateTransform|animateMotion|set)\b[^>]*>', text
    ):
        tag = m.group(0)
        begin = 0
        dur = 0
        b = re.search(r'begin="([^"]+)"', tag)
        d = re.search(r'dur="([^"]+)"', tag)
        if b:
            begin = _parse_time(b.group(1))
        if d:
            dur = _parse_time(d.group(1))
        max_end = max(max_end, begin + dur)
        # Track the longest repeating animation cycle
        if 'repeatCount="indefinite"' in tag and dur > 0:
            max_repeating_dur = max(max_repeating_dur, dur)
    # Also check standalone dur= in case some are on outer elements
    for d in re.findall(r'dur="([^"]+)"', text):
        max_end = max(max_end, _parse_time(d))
    result = max_end or 2.0
    # For indefinitely repeating animations, one full cycle is enough
    if max_repeating_dur > 0:
        result = max(result, max_repeating_dur)
    return result


def _svg_loop_start(svg_path):
    """Return the time (in seconds) after which all one-shot animations are done.

    For SVGs that mix fill="freeze" (one-shot) and repeatCount="indefinite"
    animations, this returns when the intro is finished so the GIF can loop
    from just the repeating part.  Returns 0 if there are no one-shot intros.
    """
    text = Path(svg_path).read_text(errors="ignore")
    max_oneshot_end = 0
    has_repeat = False
    max_repeating_dur = 0
    for m in re.finditer(
        r'<(?:animate|animateTransform|animateMotion|set)\b[^>]*>', text
    ):
        tag = m.group(0)
        begin = 0
        dur = 0
        b = re.search(r'begin="([^"]+)"', tag)
        d = re.search(r'dur="([^"]+)"', tag)
        if b:
            begin = _parse_time(b.group(1))
        if d:
            dur = _parse_time(d.group(1))
        if 'repeatCount="indefinite"' in tag:
            has_repeat = True
            if dur > 0:
                max_repeating_dur = max(max_repeating_dur, dur)
        elif 'fill="freeze"' in tag:
            max_oneshot_end = max(max_oneshot_end, begin + dur)
    # Only return a loop start if there are both one-shot and repeating anims
    if has_repeat and max_oneshot_end > 0:
        return max_oneshot_end, max_repeating_dur
    return 0, 0


def svg_to_png(svg_path, size=256, width=None, height=None):
    """Render SVG as a static PNG (last animation frame)."""
    svg_path = Path(svg_path)
    w = width or size
    h = height or size
    png_name = f"{svg_path.stem}_{w}x{h}.png"
    png_path = os.path.join(_tmp, png_name)
    if not os.path.exists(png_path):
        page = _browser.new_page(
            viewport={"width": w, "height": h},
            device_scale_factor=2,
        )
        page.goto(f"file://{svg_path.resolve()}")
        page.evaluate("""() => {
            const svg = document.querySelector('svg');
            if (svg && svg.pauseAnimations) {
                svg.pauseAnimations();
                svg.setCurrentTime(10);
            }
        }""")
        page.screenshot(path=png_path, omit_background=True)
        page.close()
    return png_path


def _capture_frame(page, t):
    """Set SVG to time t, wait for repaint, and return an RGBA PIL Image."""
    page.evaluate(f"""() => {{
        const svg = document.querySelector('svg');
        if (!svg) return;
        if (svg.unpauseAnimations) svg.unpauseAnimations();
        if (svg.setCurrentTime) svg.setCurrentTime({t});
        if (svg.pauseAnimations) svg.pauseAnimations();
    }}""")
    # Double-rAF ensures the browser has fully repainted
    page.evaluate(
        "() => new Promise(r => requestAnimationFrame(() => requestAnimationFrame(r)))"
    )
    buf = page.screenshot(omit_background=True)
    return Image.open(BytesIO(buf)).convert("RGBA")


def svg_to_gif(svg_path, size=256, width=None, height=None, hold_end_ms=0,
               loop_only=False):
    """Capture SVG animation frames and stitch into a looping animated GIF.

    Structure: animation frames → hold final state → loop.
    This avoids a blank flash on loop since the last & first visible
    frames are the fully-rendered end state.
    hold_end_ms: milliseconds to hold the final frame before looping.
    loop_only: if True and the SVG has a one-shot intro followed by
               repeating animations, capture only the repeating part.
    """
    svg_path = Path(svg_path)
    w = width or size
    h = height or size
    hold_suffix = f"_hold{hold_end_ms}" if hold_end_ms else ""
    loop_suffix = "_loop" if loop_only else ""
    gif_name = f"{svg_path.stem}_{w}x{h}{hold_suffix}{loop_suffix}.gif"
    gif_path = os.path.join(_tmp, gif_name)
    if os.path.exists(gif_path):
        return gif_path

    dur = _svg_animation_duration(svg_path)
    loop_start, repeating_dur = _svg_loop_start(svg_path)
    page = _browser.new_page(
        viewport={"width": w, "height": h},
        device_scale_factor=2,
    )
    page.goto(f"file://{svg_path.resolve()}")
    page.wait_for_timeout(200)

    if loop_only and loop_start > 0 and repeating_dur > 0:
        # Capture only the repeating part (skip one-shot intro)
        loop_dur = repeating_dur
        frame_count = min(GIF_MAX_FRAMES, max(10, int(loop_dur * 1000 / GIF_FRAME_DELAY_MS)))
        anim_frames = []
        for i in range(frame_count):
            t = loop_start + (i / frame_count) * loop_dur
            anim_frames.append(_capture_frame(page, t))
    else:
        # Capture the full animation from t=0 → t=dur
        frame_count = min(GIF_MAX_FRAMES, max(10, int(dur * 1000 / GIF_FRAME_DELAY_MS)))
        anim_frames = []
        for i in range(frame_count):
            t = (i / frame_count) * dur
            anim_frames.append(_capture_frame(page, t))

    # 2. Use the last animation frame as the hold frame.
    #    Capturing at exactly t=dur can cause a blank frame in some SVGs
    #    because the animation timeline resets, so we use the penultimate
    #    point (t ≈ 0.97*dur) which is the last captured frame.
    hold_frame = anim_frames[-1]

    page.close()

    # 3. Build the GIF: just the animation frames, looping continuously.
    #    Prepend one copy of the end state so the poster/first frame
    #    shows the fully-rendered image (not the blank t=0 state).
    hold_count = max(1, hold_end_ms // GIF_FRAME_DELAY_MS) if hold_end_ms else 1
    frames = [hold_frame.copy() for _ in range(hold_count)] + anim_frames

    # Save RGBA frames as transparent GIF.
    # Pillow handles RGBA→palette conversion with transparency automatically.
    frames[0].save(
        gif_path,
        save_all=True,
        append_images=frames[1:],
        duration=GIF_FRAME_DELAY_MS,
        loop=0,
        disposal=2,
    )
    return gif_path


def svg_to_image(svg_path, size=256, width=None, height=None, hold_end_ms=0):
    """Return an animated GIF if the SVG has animations, otherwise a static PNG."""
    if _svg_has_animation(svg_path):
        return svg_to_gif(svg_path, size, width, height, hold_end_ms=hold_end_ms)
    return svg_to_png(svg_path, size, width, height)


def add_icon(slide, svg_path, left, top, size=Inches(0.6)):
    """Add an SVG icon (as animated GIF or static PNG) to a slide."""
    img = svg_to_image(svg_path)
    return slide.shapes.add_picture(img, left, top, size, size)

# ── Theme colours ──────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)   # --bg: #0f172a
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)   # --accent: #38bdf8
TEXT      = RGBColor(0xF1, 0xF5, 0xF9)   # --text: #f1f5f9
MUTED     = RGBColor(0xB0, 0xBE, 0xC5)   # --text-muted
CARD_BG   = RGBColor(0x33, 0x41, 0x55)   # --card-bg
CARD_BDR  = RGBColor(0x47, 0x55, 0x69)   # --card-border
WARNING   = RGBColor(0xFB, 0xBF, 0x24)   # amber warning
COMPARE_A = RGBColor(0x1E, 0x29, 0x3B)   # darker bg for compare col
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xC0, 0x00, 0x00)   # red for quiz subtitle
GREEN_TXT = RGBColor(0xF7, 0xFF, 0xF6)   # slightly green-tinted white

# ── Additional asset paths ─────────────────────────────────────
FINAL_ASSETS = ASSETS

SLIDE_W = Emu(9144000)    # standard 16:9
SLIDE_H = Emu(5143500)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ────────────────────────────────────────────────────

def set_slide_bg(slide, colour=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_bg_image(slide, img_path):
    """Add a full-slide background image."""
    slide.shapes.add_picture(
        str(img_path), Emu(0), Emu(0), SLIDE_W, SLIDE_H)


def add_text(slide, text, left, top, width, height, *,
             font_size=18, bold=False, colour=TEXT, alignment=PP_ALIGN.LEFT,
             font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = colour
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = 1.2
    return txBox


def add_label(slide, text, top=Inches(0.35)):
    add_text(slide, text.upper(),
             Inches(0.6), top, Inches(3), Inches(0.3),
             font_size=9, bold=True, colour=ACCENT)


def add_title(slide, plain, highlight, top=Inches(0.65)):
    txBox = slide.shapes.add_textbox(Inches(0.6), top, Inches(8.8), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = plain
    run1.font.size = Pt(27)
    run1.font.bold = True
    run1.font.color.rgb = TEXT
    run1.font.name = "Segoe UI"
    p.line_spacing = 1.2
    run2 = p.add_run()
    run2.text = highlight
    run2.font.size = Pt(27)
    run2.font.bold = True
    run2.font.color.rgb = ACCENT
    run2.font.name = "Segoe UI"
    return txBox


def add_subtitle(slide, text, top=Inches(1.35), alignment=PP_ALIGN.LEFT):
    add_text(slide, text,
             Inches(0.6), top, Inches(8.8), Inches(0.4),
             font_size=15, colour=MUTED, alignment=alignment)


def add_section_divider(title_text):
    """Add a bold accent-coloured section divider slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_text(sl, title_text,
             Inches(0.5), Inches(1.5), Inches(9), Inches(2),
             font_size=40, bold=True, colour=ACCENT,
             alignment=PP_ALIGN.CENTER)
    return sl


def add_card(slide, left, top, width, height, texts, *,
             bg=CARD_BG, border=CARD_BDR):
    """texts: list of (text, font_size, bold, colour) tuples."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    for i, (txt, sz, bold, clr) in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = clr
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
    return shape


def add_bullet_list(slide, items, left, top, width, height, *,
                    font_size=15, colour=TEXT, spacing=Pt(6)):
    """items: list of (bold_text, desc_text) or plain strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.line_spacing = 1.2
        p.font.name = "Segoe UI"
        if isinstance(item, tuple):
            bold_part, desc_part = item
            r1 = p.add_run()
            r1.text = bold_part
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = TEXT
            r1.font.name = "Segoe UI"
            if desc_part:
                r2 = p.add_run()
                r2.text = f"  —  {desc_part}"
                r2.font.size = Pt(font_size - 2)
                r2.font.bold = False
                r2.font.color.rgb = MUTED
                r2.font.name = "Segoe UI"
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = colour
    return txBox


def add_icon_rows(slide, rows, *, left=Inches(0.75), top=Inches(1.8),
                  row_h=Inches(0.67), icon_size=Inches(0.41),
                  text_width=Inches(7.5), font_size=15, desc_size=None,
                  text_h=None):
    """Place icon+text pairs at exact Y positions so they stay aligned.

    rows: list of (svg_path, bold_text, desc_text)
    """
    if desc_size is None:
        desc_size = font_size - 2
    if text_h is None:
        text_h = row_h
    text_left = left + icon_size + Inches(0.2)
    for i, (svg, bold_txt, desc_txt) in enumerate(rows):
        y = top + i * row_h
        icon_y = y + (row_h - icon_size) / 2
        add_icon(slide, svg, left, icon_y, icon_size)
        txBox = slide.shapes.add_textbox(text_left, y, text_width, text_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = bold_txt
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = TEXT
        r1.font.name = "Segoe UI"
        p.line_spacing = 1.2
        if desc_txt:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            p2.line_spacing = 1.2
            r2 = p2.add_run()
            r2.text = desc_txt
            r2.font.size = Pt(desc_size)
            r2.font.bold = False
            r2.font.color.rgb = MUTED
            r2.font.name = "Segoe UI"


# ══════════════════════════════════════════════════════════════
# SLIDES — ordered to match AI4Students_April_2026_Final.pptx
# ══════════════════════════════════════════════════════════════

# ── Slide 1: Title ─────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(sl)
add_bg_image(sl, FINAL_ASSETS / "bg-title.png")
add_text(sl, "AI & the Future",
         Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
         font_size=50, bold=True, colour=ACCENT,
         alignment=PP_ALIGN.CENTER)
add_text(sl, "Artificial Intelligence · Generative AI · Agents · Vibe Coding",
         Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
         font_size=22, colour=MUTED, alignment=PP_ALIGN.CENTER)


# ── Slide 2: Agenda ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Overview")
add_title(sl, "Session ", "Overview")
agenda = [
    "1.  Welcome & Introduction",
    "2.  Ice Breaker & Quiz",
    "3.  How we got here — from Machine Learning to Agents",
    "4.  Generative AI & using it responsibly",
    "5.  Vibe Coding — what & why",
    "6.  Interactive Q&A & demos",
    "7.  Hands‑on Vibe Coding session",
    "8.  Showcase & wrap‑up",
]
add_bullet_list(sl, agenda,
                Inches(1.1), Inches(1.6), Inches(7.5), Inches(3.8),
                font_size=15, colour=TEXT)


# ── Slide 3: Team Members ─────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Your Hosts")
add_title(sl, "Team ", "Members")

team = [
    ("RP", "Ranjan Patalasingh", "Principal AI Engineer"),
    ("MM", "Muhammed Miah", "Senior AI Engineer"),
    ("FH", "Farhad Hussain", "Senior AI Engineer"),
    ("KC", "Kelvin Chan", "Data Scientist"),
    ("CK", "Chaitanya Katukuri", "Engineering Lead, ML"),
    ("KK", "Krishna Konduru", "Data Engineer"),
]
# SVG avatars for all members
team_avatars = {
    "RP": FINAL_ASSETS / "avatar-rp.svg",
    "MM": FINAL_ASSETS / "avatar-mm.svg",
    "FH": FINAL_ASSETS / "avatar-fh.svg",
    "KC": FINAL_ASSETS / "avatar-kc.svg",
    "CK": FINAL_ASSETS / "avatar-ck.svg",
    "KK": FINAL_ASSETS / "avatar-kk.svg",
}
card_w = Inches(2.6)
card_h = Inches(1.35)
gap = Inches(0.3)
total_team_w = 3 * card_w + 2 * gap
start_x = (SLIDE_W - total_team_w) / 2
for i, (initials, name, role) in enumerate(team):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap)
    y = Inches(1.7) + row * (card_h + gap)
    card_shape = add_card(sl, x, y, card_w, card_h, [
        (initials, 18, True, ACCENT),
        ("", 6, False, MUTED),
        (name, 12, False, TEXT),
        (role, 9.75, False, MUTED),
    ])
    card_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Add SVG avatar — placed inside the card near the top
    photo_sz = Inches(0.67)
    photo_x = x + (card_w - photo_sz) / 2
    photo_y = y + Inches(0.06)
    if initials in team_avatars:
        avatar_png = svg_to_png(team_avatars[initials], 200, 200)
        if avatar_png:
            sl.shapes.add_picture(str(avatar_png), photo_x, photo_y, photo_sz, photo_sz)


# ── Slide 4: Section Divider — Ice Breaker & Quiz ──────────────
add_section_divider("Ice Breaker & Quiz")


# ── Slide 5: Ice Breaker ──────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Warm Up")
add_title(sl, "Let's Break the ", "Ice!")
add_subtitle(sl, "Quick‑fire round — shout your answer!")
add_icon_rows(sl, [
    (ASSETS / "icon-robot.svg",
     "What AI tool have you used in the last week?", ""),
    (ASSETS / "icon-lightbulb.svg",
     "Name one thing you think AI cannot do (yet).", ""),
    (VIBE_ASSETS / "icon-rocket.svg",
     "If you could build any app with AI, what would it be?", ""),
    (VIBE_ASSETS / "icon-target.svg",
     "What do you most want to learn today?", ""),
], left=Inches(0.9), top=Inches(1.95), row_h=Inches(0.75),
   icon_size=Inches(0.41), font_size=16.5)


# ── Slide 6: Quiz ─────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Challenge")
add_title(sl, "", "Quiz")
# Quiz image
try:
    quiz_img = FINAL_ASSETS / "quiz-image.png"
    if quiz_img.exists():
        sl.shapes.add_picture(
            str(quiz_img),
            Inches(1.15), Inches(1.28), Inches(7.7), Inches(4.0))
except Exception:
    pass
add_text(sl, "Fast facts — first answer wins",
         Inches(0.5), Inches(3.0), Inches(9), Inches(0.5),
         font_size=28, bold=True, colour=RED, alignment=PP_ALIGN.CENTER)


# ── Slide 7: Section Divider — ML, GenAI & Agents ─────────────
add_section_divider("Machine Learning, Generative AI & Agents")


# ── Slide 8: The Hook ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "The Big Idea")
txBox = sl.shapes.add_textbox(
    Inches(0.6), Inches(1.1), Inches(8.8), Inches(1.1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "What if you could build a website\njust by describing it?"
r.font.size = Pt(27)
r.font.bold = True
r.font.color.rgb = ACCENT
r.font.name = "Segoe UI"
add_text(sl, "No tutorials. No experience needed. Just your ideas.",
         Inches(0.5), Inches(2.6), Inches(9), Inches(0.4),
         font_size=16.5, colour=MUTED, alignment=PP_ALIGN.CENTER)
# hero-flow graphic
try:
    hero_w = Inches(6.77)
    hero_h = Inches(1.88)
    png = svg_to_image(VIBE_ASSETS / "hero-flow.svg", width=720, height=200, hold_end_ms=500)
    sl.shapes.add_picture(png,
                          Inches(1.65), Inches(3.15), hero_w, hero_h)
except Exception:
    pass


# ── Slide 9: How We Got Here ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "The Journey")
add_title(sl, "How We Got Here — From Learning to ", "Doing")
add_subtitle(sl, "AI evolved in three big leaps — each one unlocked the next.")

roles = [
    ("The Analyst", "Machine Learning & Deep Learning · ~2010s",
     "Learns from data, recognises patterns, makes predictions. Classifies what already exists.",
     "Spotify recommendations, FaceID, spam filters."),
    ("The Creator", "LLMs & Generative AI · ~2018–2022",
     "Understands language, then starts generating new things — text, images, code, music.",
     "ChatGPT writing an essay,\nMidjourney generating artwork."),
    ("The Doer", "Agents · ~2024",
     "AI that can act autonomously —\nbrowses the web, runs code, plans multi-step tasks.",
     "GitHub Copilot building features for you."),
]
role_icons = [
    ASSETS / "icon-analyst.svg",
    ASSETS / "icon-creator.svg",
    ASSETS / "icon-doer.svg",
]
card_w = Inches(2.8)
card_gap = Inches(0.22)
card_top = Inches(1.95)
icon_sz = Inches(0.45)
total_w = 3 * card_w + 2 * card_gap
card_left = (SLIDE_W - total_w) / 2
for i, (title, era, desc, eg) in enumerate(roles):
    x = card_left + i * (card_w + card_gap)
    card_shape = add_card(sl, x, card_top, card_w, Inches(2.4), [
        (title, 15, True, ACCENT),
        (era, 8.25, False, MUTED),
        (desc, 9.75, False, TEXT),
        ("", 6, False, MUTED),
        (eg, 8.25, False, MUTED),
    ])
    card_shape.text_frame.margin_top = Inches(0.6)
    add_icon(sl, role_icons[i],
             x + (card_w - icon_sz) / 2,
             card_top + Inches(0.18),
             icon_sz)

# Payoff banner — no card/box, just icon + text
add_icon(sl, VIBE_ASSETS / "icon-rocket.svg", Inches(3.11), Inches(4.76), Inches(0.45))
add_text(sl, "The Payoff — Vibe Coding (~2025)",
         Inches(3.48), Inches(4.81), Inches(3.38), Inches(0.35),
         font_size=13.5, bold=True, colour=GREEN_TXT)


# ── Slide 10: ML, Deep Learning and LLM ───────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_text(sl, "2010s",
         Inches(0.6), Inches(0.35), Inches(3), Inches(0.3),
         font_size=9, bold=True, colour=ACCENT)
add_title(sl, "ML, Deep Learning and ", "LLM")

# Each definition is positioned next to its diagram
ml_defs = [
    ("Machine learning is a subset of artificial intelligence that enables systems to learn from data without being explicitly programmed.",
     Inches(2.0), Inches(1.44), Inches(3.2), Inches(0.63), PP_ALIGN.RIGHT),
    ("Deep learning is a subset of machine learning that uses neural networks to analyze data.",
     Inches(5.2), Inches(2.92), Inches(3.13), Inches(0.45), PP_ALIGN.LEFT),
    ("Large language models are a type of Deep Learning model specifically designed for natural language processing tasks.",
     Inches(2.43), Inches(4.23), Inches(2.91), Inches(0.63), PP_ALIGN.RIGHT),
]
for defn, dl, dt, dw, dh, align in ml_defs:
    add_text(sl, defn, dl, dt, dw, dh,
             font_size=10.5, bold=False, colour=TEXT, alignment=align)

# Diagrams
diagrams = [
    ("diagram-ml.png", Inches(5.35), Inches(1.17), Inches(4.24), Inches(1.54)),
    ("diagram-dl.png", Inches(0.42), Inches(2.38), Inches(4.57), Inches(1.56)),
    ("diagram-llm.png", Inches(5.35), Inches(3.67), Inches(4.14), Inches(1.87)),
]
for fname, dl, dt, dw, dh in diagrams:
    img_path = FINAL_ASSETS / fname
    if img_path.exists():
        sl.shapes.add_picture(str(img_path), dl, dt, dw, dh)


# ── Slide 11: Generative AI ────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "In the Wild")
add_title(sl, "Generative AI — Where You ", "See It Today")

# How it works
add_text(sl, "How it works", Inches(0.6), Inches(1.43), Inches(3), Inches(0.3),
         font_size=12, bold=True, colour=ACCENT)
add_icon_rows(sl, [
    (VIBE_ASSETS / "icon-describe.svg",
     "You write a prompt", "Describe what you want in plain language"),
    (ASSETS / "icon-flow.svg",
     "The model processes it", "Breaks your words into tokens and predicts what comes next"),
    (ASSETS / "icon-genai.svg",
     "You get the output", "Text, an image, code, audio, or video"),
], left=Inches(0.6), top=Inches(1.95), row_h=Inches(0.75),
   icon_size=Inches(0.375), text_width=Inches(3.9), font_size=10.5, desc_size=9,
   text_h=Inches(0.46))

# Tools grid
add_text(sl, "Tools you might already know", Inches(5.25), Inches(1.43), Inches(4.12), Inches(0.3),
         font_size=12, bold=True, colour=ACCENT)
tool_icons = [
    ASSETS / "icon-chat.svg",
    ASSETS / "icon-image-video.svg",
    ASSETS / "icon-code.svg",
    ASSETS / "icon-audio.svg",
]
tools = [
    ("Text & Chat", "ChatGPT, Claude, Gemini, Grok"),
    ("Images & Video", "GPT Image, Midjourney, Google Veo"),
    ("Code Generation", "GitHub Copilot, Cursor, Claude Code"),
    ("Audio & Music", "ElevenLabs, Suno, Udio"),
]
tw = Inches(2.02)
th = Inches(0.90)
tool_icon_sz = Inches(0.3)
tool_txt_h = Inches(0.63)
for i, (name, desc) in enumerate(tools):
    col = i % 2
    row = i // 2
    x = Inches(5.25) + col * (tw + Inches(0.15))
    y = Inches(1.88) + row * (th + Inches(0.14))
    txt_x = x + Inches(0.11) + tool_icon_sz + Inches(0.08)
    txt_w = Inches(1.43)
    add_card(sl, x, y, tw, th, [], bg=CARD_BG, border=CARD_BDR)
    add_icon(sl, tool_icons[i], x + Inches(0.11), y + Inches(0.29), tool_icon_sz)
    txBox = sl.shapes.add_textbox(txt_x, y + Inches(0.13), txt_w, tool_txt_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(8.25)
    p2.font.color.rgb = MUTED
    p2.font.name = "Segoe UI"

# Warning banner
add_card(sl, Inches(2.0), Inches(4.47), Inches(6.0), Inches(0.75), [], bg=CARD_BG, border=CARD_BDR)
add_icon(sl, ASSETS / "icon-warning.svg", Inches(2.2), Inches(4.66), Inches(0.375))
txBox = sl.shapes.add_textbox(Inches(2.77), Inches(4.50), Inches(5.12), Inches(0.69))
tf = txBox.text_frame
tf.word_wrap = True
tf.margin_left = 0
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "⚠  Hallucination warning"
p.font.size = Pt(10.5)
p.font.bold = True
p.font.color.rgb = WARNING
p.font.name = "Segoe UI"
p.line_spacing = 1.2
p2 = tf.add_paragraph()
p2.text = "Generative AI can confidently produce wrong answers. It sounds right, reads right, but it's made up. Always verify."
p2.font.size = Pt(9.75)
p2.font.color.rgb = TEXT
p2.font.name = "Segoe UI"
p2.line_spacing = 1.2


# ── Slide 12: Humans vs AI Benchmarks ──────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Benchmarks")
# Title as two separate paragraphs to match final
txBox = sl.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.9))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "Humans VS"
r1.font.size = Pt(27)
r1.font.bold = True
r1.font.color.rgb = TEXT
r1.font.name = "Segoe UI"
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "Artificial Intelligence"
r2.font.size = Pt(27)
r2.font.bold = True
r2.font.color.rgb = TEXT
r2.font.name = "Segoe UI"

# Single text block with all benchmark lines
benchmark_lines = [
    "This infographic tracks how AI performance has evolved relative to human baseline (100%) across several technical domains from 2012 to 2024.",
    "Visual Reasoning & Medium-Level Reading Comprehension",
    "English Language Understanding",
    "Multitask Language Understanding",
    "Competition-Level Mathematics",
    "PhD-Level Science Questions",
    "Multimodal Understanding & Reasoning",
]
txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(4.76), Inches(2.99))
tf = txBox.text_frame
tf.word_wrap = True
for j, line in enumerate(benchmark_lines):
    if j == 0:
        p = tf.paragraphs[0]
    elif j == 1:
        # Empty line between intro sentence and bullet list
        spacer = tf.add_paragraph()
        spacer.space_before = Pt(0)
        spacer.space_after = Pt(0)
        sr = spacer.add_run()
        sr.text = ""
        sr.font.size = Pt(6)
        p = tf.add_paragraph()
    else:
        p = tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.size = Pt(12)
    r.font.color.rgb = GREEN_TXT
    r.font.name = "Segoe UI"
    p.line_spacing = 1.2
    if j > 0:
        p.level = 1
        pPr = p._p.get_or_add_pPr()
        marL = str(int(Inches(0.3)))
        indent = str(int(-Inches(0.15)))
        pPr.set('marL', marL)
        pPr.set('indent', indent)
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': 'F7FFF6'})
        buClr.append(srgbClr)
        pPr.append(buClr)
        buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})
        pPr.append(buSzPct)
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        pPr.append(buChar)
        p.space_before = Pt(4)

# Benchmark chart image
try:
    chart_img = FINAL_ASSETS / "benchmark-chart.png"
    if chart_img.exists():
        sl.shapes.add_picture(
            str(chart_img),
            Inches(5.47), Inches(0.34), Inches(4.3), Inches(5.04))
except Exception:
    pass


# ── Slide 13: Section Divider — Vibe Coding ────────────────────
add_section_divider("Vibe Coding")


# ── Slide 14: What is Vibe Coding? ─────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "The Idea")
add_title(sl, "What is ", "Vibe Coding?")
add_icon_rows(sl, [
    (VIBE_ASSETS / "icon-describe.svg",
     "You describe", "Tell AI what you want in plain English"),
    (VIBE_ASSETS / "icon-ai.svg",
     "AI builds", "GitHub Copilot writes the code for you"),
    (VIBE_ASSETS / "icon-browser.svg",
     "You see results", "Open the browser and watch it come to life"),
], left=Inches(0.75), top=Inches(1.7), row_h=Inches(0.9),
   icon_size=Inches(0.45), text_width=Inches(7.5), font_size=16.5, desc_size=15)


# ── Slide 15: Why Vibe Coding Matters ──────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Why It Matters")
add_title(sl, "Why Vibe Coding ", "Matters")
add_icon_rows(sl, [
    (ASSETS / "icon-nosyntax.svg",
     "No syntax memorisation needed", "Describe ideas in plain English"),
    (ASSETS / "icon-iterate.svg",
     "Iterate in seconds, not hours", "See results in your browser almost immediately"),
    (ASSETS / "icon-learn.svg",
     "Learn by doing", "Read & edit AI-generated code to understand how it works"),
    (ASSETS / "icon-prototype.svg",
     "Ideal for prototypes & side projects", "Go from idea to working thing fast"),
], left=Inches(0.6), top=Inches(1.7), row_h=Inches(0.67),
   icon_size=Inches(0.41), text_width=Inches(4.7), font_size=13.5, desc_size=12)

# Comparison grid
cw = Inches(2.4)
ch = Inches(1.58)
comp_x = Inches(6.2)
add_card(sl, comp_x, Inches(1.35), cw, ch, [
    ("Traditional coding", 12, True, TEXT),
    ("Write every line manually", 9.75, False, MUTED),
    ("Debug from scratch", 9.75, False, MUTED),
    ("Steep learning curve", 9.75, False, MUTED),
], bg=COMPARE_A)
add_card(sl, comp_x, Inches(3.08), cw, ch, [
    ("Vibe coding", 12, True, ACCENT),
    ("Describe → review → refine", 9.75, False, TEXT),
    ("AI helps find & fix issues", 9.75, False, TEXT),
    ("Low barrier to entry", 9.75, False, TEXT),
])


# ── Slide 16: What's Possible ─────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "What's Possible")
add_title(sl, "People have already built ", "all of this")
PROJ_ROOT = HERE.parent
projects = [
    ("Pixel Art Editor", PROJ_ROOT / "pixel-art-editor" / "assets" / "thumbnail.png"),
    ("Snake Game", PROJ_ROOT / "snake-with-powerups" / "assets" / "thumbnail.png"),
    ("Beat Maker", PROJ_ROOT / "beat-maker" / "assets" / "thumbnail.png"),
    ("Personality Quiz", PROJ_ROOT / "personality-quiz" / "assets" / "thumbnail.png"),
    ("Dungeon Crawler", PROJ_ROOT / "dungeon-crawler" / "assets" / "thumbnail.png"),
    ("Weather Dashboard", PROJ_ROOT / "weather-dashboard" / "assets" / "thumbnail.png"),
    ("Endless Runner", PROJ_ROOT / "endless-runner" / "assets" / "thumbnail.png"),
    ("Memory Match", PROJ_ROOT / "memory-card-match" / "assets" / "thumbnail.png"),
]
img_w = Inches(1.99)
img_h = Inches(1.09)
lbl_w = Inches(2.14)
lbl_h = Inches(0.27)
col_lefts = [Inches(0.57), Inches(2.86), Inches(5.15), Inches(7.44)]
txt_lefts = [Inches(0.50), Inches(2.79), Inches(5.07), Inches(7.36)]
row_img_tops = [Inches(1.56), Inches(3.21)]
row_txt_tops = [Inches(2.66), Inches(4.31)]
for i, (name, thumb) in enumerate(projects):
    col = i % 4
    row = i // 4
    pic = sl.shapes.add_picture(str(thumb), col_lefts[col], row_img_tops[row],
                                img_w, img_h)
    prstGeom = pic._element.spPr.find(qn('a:prstGeom'))
    prstGeom.set('prst', 'roundRect')
    avLst = prstGeom.find(qn('a:avLst'))
    gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val 5000'})
    avLst.append(gd)
    add_text(sl, name, txt_lefts[col], row_txt_tops[row], lbl_w, lbl_h,
             font_size=9.75, bold=True, colour=TEXT, alignment=PP_ALIGN.CENTER)
add_subtitle(sl, "All built by describing ideas in plain English — no coding experience needed.",
             top=Inches(4.99), alignment=PP_ALIGN.CENTER)


# ── Slide 17: Q&A ─────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Ask Us Anything")
add_title(sl, "Interactive Q&A & ", "Live Demo")
add_text(sl, "Got questions? Curious about something?\nNow's the time — ask us anything!",
         Inches(0.5), Inches(2.4), Inches(9), Inches(0.75),
         font_size=18, colour=MUTED, alignment=PP_ALIGN.CENTER)


# ── Slide 18: Section Divider — Hands-on Vibe Coding ───────────
add_section_divider("Hands\u2011on Vibe Coding")


# ── Slide 19: Tips ─────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Let's Go!")
add_title(sl, "Quick tips before you ", "start")
tip_icons = [
    VIBE_ASSETS / "icon-target.svg",
    VIBE_ASSETS / "icon-blocks.svg",
    VIBE_ASSETS / "icon-flask.svg",
]
tips = [
    ("🎯 Be specific",
     "\"Add a red button that plays a drum sound\" works way better than \"make it cool\""),
    ("🧱 Build step by step",
     "Start simple, then keep adding. Don't try to describe everything at once."),
    ("🧪 Experiment!",
     "Nothing can break permanently. Just try things and see what happens."),
]
tw = Inches(2.8)
card_top = Inches(2.1)
tip_icon_sz = Inches(0.45)
tip_gap = Inches(0.22)
total_tip_w = 3 * tw + 2 * tip_gap
tip_left = (SLIDE_W - total_tip_w) / 2
for i, (title, desc) in enumerate(tips):
    x = tip_left + i * (tw + tip_gap)
    card_shape = add_card(sl, x, card_top, tw, Inches(2.25), [
        (title, 16.5, True, TEXT),
        ("", 6, False, MUTED),
        (desc, 11.25, False, MUTED),
    ])
    card_shape.text_frame.margin_top = Inches(0.63)
    add_icon(sl, tip_icons[i],
             x + (tw - tip_icon_sz) / 2,
             card_top + Inches(0.26),
             tip_icon_sz)


# ── Slide 20: Hands-On ────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Your Turn")
add_title(sl, "Hands\u2011On: Build Something ", "Today!")

add_icon_rows(sl, [
    (VIBE_ASSETS / "icon-idea.svg",
     "1. Pick your idea", "Choose a starter or dream up your own"),
    (VIBE_ASSETS / "icon-describe.svg",
     "2. Build with AI", "Describe what you want in plain English"),
    (VIBE_ASSETS / "icon-browser.svg",
     "3. Test & fix", "See results, tweak, iterate"),
    (ASSETS / "icon-showcase.svg",
     "4. Demo to the group", "Show off what you built!"),
], left=Inches(0.6), top=Inches(1.65), row_h=Inches(0.75),
   icon_size=Inches(0.375), text_width=Inches(3.9), font_size=13.5, desc_size=12)

# Project ideas
add_text(sl, "Project ideas (or choose your own!)",
         Inches(5.25), Inches(1.5), Inches(4.1), Inches(0.3),
         font_size=12, bold=True, colour=ACCENT)
ideas = [
    ("🧠 Quiz generator", "Build a web quiz on any topic you love"),
    ("🍳 Recipe recommender", "Type ingredients, get a recipe"),
    ("📚 Study planner", "AI that schedules A-Level revision"),
    ("🗣️ Chatbot persona", "A bot that talks like a historical figure"),
    ("🎮 Mini game", "Text-based adventure or number guessing game"),
    ("🌍 Climate dashboard", "Visualise climate data with charts and facts"),
]
iw = Inches(2.0)
ih = Inches(0.975)
for i, (title, desc) in enumerate(ideas):
    col = i % 2
    row = i // 2
    x = Inches(5.25) + col * (iw + Inches(0.1))
    y = Inches(1.95) + row * (ih + Inches(0.1))
    add_card(sl, x, y, iw, ih, [
        (title, 9.75, True, TEXT),
        (desc, 8.25, False, MUTED),
    ])


# ── Slide 21: Showcase & Wrap-Up ──────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Finale")
add_title(sl, "Showcase & ", "Wrap\u2011Up")
add_subtitle(sl, "Each group: 60–120 seconds to demo what you built.")

add_text(sl, "Key takeaways", Inches(0.6), Inches(1.95), Inches(3), Inches(0.3),
         font_size=12, bold=True, colour=ACCENT)
add_icon_rows(sl, [
    (ASSETS / "icon-analyst.svg",
     "AI is a tool", "The human with the best questions wins"),
    (ASSETS / "icon-creator.svg",
     "Analyst → Creator → Doer",
     "The Analyst recognises → the Creator generates → the Doer acts — each built on the last"),
    (VIBE_ASSETS / "icon-rocket.svg",
     "Vibe coding", "Lets you ship ideas without years of training"),
], left=Inches(0.75), top=Inches(2.4), row_h=Inches(0.75),
   icon_size=Inches(0.45), text_width=Inches(7.5), font_size=15, desc_size=13.5)


# ── Slide 22: Section Divider — Warning & Risks ───────────────
add_section_divider("WARNING & RISKS")


# ── Slide 23: Pitfalls of Vibe Coding ─────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Warning")
add_title(sl, "Pitfalls of ", "Vibe Coding")
add_icon_rows(sl, [
    (ASSETS / "icon-lock.svg",
     "Security Risks",
     "Vibe-coding can lead to hidden vulnerabilities, exposing systems to potential threats and compromising security."),
    (ASSETS / "icon-tools.svg",
     "Maintainability Challenges",
     "Incremental prompting typically results in poorly structured code, making debugging, extending, or refactoring difficult."),
    (ASSETS / "icon-court.svg",
     "Compliance Issues",
     "Coding without proper understanding may bypass regulations, leading to compliance problems and legal risks."),
], left=Inches(0.6), top=Inches(1.5), row_h=Inches(1.1),
   icon_size=Inches(0.38), text_width=Inches(4.5), font_size=13.5, desc_size=12)
try:
    pitfalls_img = FINAL_ASSETS / "pitfalls-vibe.png"
    if pitfalls_img.exists():
        sl.shapes.add_picture(
            str(pitfalls_img),
            Inches(6.21), Inches(1.32), Inches(3.43), Inches(4.09))
except Exception:
    pass


# ── Slide 24: Use AI Responsibly ───────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_label(sl, "Stay Sharp")
add_title(sl, "Use AI ", "Responsibly")
add_subtitle(sl, "AI is powerful — but it's not perfect. Here are things to watch out for:")
add_icon_rows(sl, [
    (ASSETS / "icon-deepfake.svg",
     "Deepfakes", "AI can create convincing fake images, audio, and video of real people"),
    (ASSETS / "icon-bias.svg",
     "Bias", "AI learns from human data — so it can repeat unfair patterns and stereotypes"),
    (ASSETS / "icon-misinfo.svg",
     "Misinformation", "AI-generated text can sound authoritative even when it's completely wrong"),
    (ASSETS / "icon-copyright.svg",
     "Copyright", "AI may generate content that's too similar to someone else's work"),
], left=Inches(0.75), top=Inches(1.8), row_h=Inches(0.75),
   icon_size=Inches(0.45), text_width=Inches(7.5), font_size=15, desc_size=13.5)


# ── Slide 25: Thank You ───────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_bg_image(sl, FINAL_ASSETS / "bg-title.png")
add_text(sl, "Thank You",
         Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
         font_size=50, bold=True, colour=ACCENT,
         alignment=PP_ALIGN.CENTER)
add_text(sl, "Now go build something amazing! 🚀",
         Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
         font_size=24, colour=TEXT, alignment=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
out = "AI_and_the_Future.pptx"
prs.save(out)

_browser.close()
_pw.stop()

print(f"✅ Saved {out} ({len(prs.slides)} slides)")
